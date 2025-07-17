import os
from flask import Flask, request, render_template, jsonify, session, make_response
from flask_socketio import SocketIO, emit
from flask_cors import CORS
from docx2python import docx2python
import pandas as pd
import PyPDF2
from elasticsearch import Elasticsearch
from elasticsearch import ElasticsearchWarning
from datetime import datetime
import markdown
from pptx import Presentation
from bs4 import BeautifulSoup
import warnings
from docx2pdf import convert as docx_to_pdf
from docx import Document
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
import tempfile
import subprocess
import logging
import lxml.etree as ET
import pytesseract
pytesseract.pytesseract.tesseract_cmd = "/opt/homebrew/bin/tesseract"
from PIL import Image
from pdf2image import convert_from_path
import io
import glob
import urllib.parse
import zipfile
import concurrent.futures
import caching
import corruptzip

# Set up logging
logging.basicConfig(level=logging.DEBUG)
logger = logging.getLogger(__name__)

warnings.filterwarnings("ignore", category=ElasticsearchWarning)

# Initialize Elasticsearch client
es = Elasticsearch("http://localhost:9200")
if not es.ping():
    logger.error("Failed to connect to Elasticsearch. Ensure it's running or update the connection settings.")
    raise Exception("Elasticsearch connection failed")
logger.debug("Successfully connected to Elasticsearch")

# Supported file extensions
SUPPORTED_EXTENSIONS = {".docx", ".xlsx", ".pdf", ".txt", ".rtf", ".doc", ".md", ".csv", ".pptx", ".xml", ".png", ".jpeg", ".jpg"}

# Initialize Flask app and SocketIO
app = Flask(__name__)
app.config["SECRET_KEY"] = "a1b2c3d4e5f6g7h8i9j0k1l2m3n4o5p"
CORS(app, resources={r"/preview/*": {"origins": "*"}})
socketio = SocketIO(app, async_mode="eventlet")

if os.environ.get('FLASK_ENV') == 'development':
    app.config['public_url'] = 'http://localhost:5000'
else:
    app.config['public_url'] = os.environ.get('NGROK_URL', 'http://localhost:5000')

LIBREOFFICE_PATH = "/Applications/LibreOffice.app/Contents/MacOS/soffice"

# Custom paths for corrupt files
CORRUPT_FILES_DIR = "/Volumes/CORRUPTED/AI/corruptedfiles"
CORRUPT_FILES_ZIP = os.path.join(CORRUPT_FILES_DIR, "corrupt_files.zip")

def extract_text_from_image(image):
    try:
        logger.debug("Starting OCR on image")
        text = pytesseract.image_to_string(image, config='--psm 6')
        logger.debug(f"OCR result: {text[:50]}...")
        return text
    except Exception as e:
        logger.error(f"Error in OCR: {e}")
        return ""

def extract_text(file_path, skip_ocr=False):
    """Base extraction function, wrapped by caching and corrupt handling."""
    ext = os.path.splitext(file_path)[1].lower()
    content = ""
    
    if ext == ".docx":
        doc = docx2python(file_path)
        content = " ".join(doc.text.split())
        logger.debug(f"Extracted from {file_path} (docx): {content[:50]}...")
        if not skip_ocr:
            with zipfile.ZipFile(file_path) as docx_zip:
                image_texts = []
                for file_name in docx_zip.namelist():
                    if file_name.startswith("word/media/") and file_name.lower().endswith((".png", ".jpeg", ".jpg")):
                        try:
                            with docx_zip.open(file_name) as img_file:
                                img = Image.open(img_file)
                                text = extract_text_from_image(img)
                                image_texts.append(text)
                        except Exception as e:
                            logger.error(f"Error extracting image {file_name} from {file_path}: {e}")
                if image_texts:
                    content += " " + " ".join(image_texts)
                    logger.debug(f"Extracted image text from {file_path} (docx): {image_texts[:50]}...")
    
    elif ext == ".xlsx":
        df = pd.read_excel(file_path)
        content = df.to_string()
        logger.debug(f"Extracted from {file_path} (xlsx): {content[:50]}...")
    
    elif ext == ".pdf":
        with open(file_path, "rb") as f:
            pdf = PyPDF2.PdfReader(f)
            content = " ".join([pdf.pages[i].extract_text() or "" for i in range(len(pdf.pages))])
            logger.debug(f"Extracted from {file_path} (pdf): {content[:50]}...")
        if not skip_ocr:
            try:
                images = convert_from_path(file_path)
                image_texts = []
                for img in images:
                    text = extract_text_from_image(img)
                    image_texts.append(text)
                if image_texts:
                    content += " " + " ".join(image_texts)
                    logger.debug(f"Extracted image text from {file_path} (pdf): {image_texts[:50]}...")
            except Exception as e:
                logger.error(f"Error extracting images from {file_path} (pdf): {e}")
    
    elif ext == ".md":
        with open(file_path, "r", encoding="utf-8") as f:
            md_content = f.read()
            html = markdown.markdown(md_content)
            soup = BeautifulSoup(html, "html.parser")
            content = soup.get_text(separator=" ")
            logger.debug(f"Extracted from {file_path} (md): {content[:50]}...")
    
    elif ext == ".csv":
        df = pd.read_csv(file_path)
        content = df.to_string()
        logger.debug(f"Extracted from {file_path} (csv): {content[:50]}...")
    
    elif ext == ".pptx":
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        content = " ".join(text)
        logger.debug(f"Extracted from {file_path} (pptx): {content[:50]}...")
    
    elif ext == ".xml":
        tree = ET.parse(file_path)
        root = tree.getroot()
        metadata = {elem.tag: elem.text for elem in root if elem.text}
        content = " ".join(root.itertext()).strip()
        logger.debug(f"Extracted from {file_path} (xml): {content[:50]}... Metadata: {metadata}")
        return content, metadata
    
    elif ext in {".png", ".jpeg", ".jpg"}:
        img = Image.open(file_path)
        if not skip_ocr:
            content = extract_text_from_image(img)
            logger.debug(f"Extracted from {file_path} (image): {content[:50]}...")
    
    return content or ""

def index_file(file_path, corrupt_dir=CORRUPT_FILES_DIR):
    content_and_status, was_cached = caching.get_cached_or_extract(file_path, lambda fp: corruptzip.extract_with_repair(fp, extract_text))
    content, metadata, is_corrupt = content_and_status
    
    if is_corrupt:
        return corruptzip.handle_corrupt_file(file_path, corrupt_dir)
    
    if not content and not metadata:
        return False
    
    doc = {
        "path": file_path,
        "content": content,
        "timestamp": datetime.now().isoformat()
    }
    if metadata:
        doc["metadata"] = metadata
    
    try:
        es.index(index="knowledge_repo", body=doc)
        logger.debug(f"Indexed: {file_path} (Cached: {was_cached})")
        return True
    except Exception as e:
        logger.error(f"Failed to index {file_path}: {e}")
        return False

def index_directory(directory, corrupt_dir=CORRUPT_FILES_DIR, corrupt_zip_path=CORRUPT_FILES_ZIP):
    if not os.path.isdir(directory):
        logger.error(f"Directory {directory} is not valid or accessible.")
        socketio.emit("index_complete", {"message": f"Error: Directory {directory} is not valid or accessible."})
        return 0, 0
    
    if es.indices.exists(index="knowledge_repo"):
        es.indices.delete(index="knowledge_repo")
    mapping = {
        "mappings": {
            "properties": {
                "path": {"type": "keyword"},
                "content": {"type": "text"},
                "timestamp": {"type": "date"},
                "metadata": {"type": "object", "enabled": True}
            }
        }
    }
    try:
        es.indices.create(index="knowledge_repo", body=mapping)
        logger.debug("Created Elasticsearch index: knowledge_repo")
    except Exception as e:
        logger.error(f"Failed to create Elasticsearch index: {e}")
        socketio.emit("index_complete", {"message": f"Error: Failed to create Elasticsearch index: {e}"})
        return 0, 0
    
    files_to_index = []
    for root, _, files in os.walk(directory):
        for file in files:
            if os.path.splitext(file)[1].lower() in SUPPORTED_EXTENSIONS and not file.startswith('.'):
                files_to_index.append(os.path.join(root, file))
    
    logger.debug(f"Found {len(files_to_index)} files to index in {directory}")
    total_files = len(files_to_index)
    indexed_files = [0]
    corrupt_files = [0]
    repaired_files = [0]
    
    def process_file(file_path):
        try:
            content_and_status, _ = caching.get_cached_or_extract(file_path, lambda fp: corruptzip.extract_with_repair(fp, extract_text))
            content, metadata, is_corrupt = content_and_status
            
            if is_corrupt:
                corrupt_files[0] += 1
                return None
            
            if not content and not metadata:
                return None
            
            if content and os.path.exists(file_path):
                repaired_files[0] += 1
            
            doc = {
                "path": file_path,
                "content": content,
                "timestamp": datetime.now().isoformat()
            }
            if metadata:
                doc["metadata"] = metadata
            logger.debug(f"Prepared document for {file_path}: {doc}")
            return doc
        except Exception as e:
            logger.error(f"Error processing {file_path}: {e}")
            return None
    
    docs_to_index = []
    with concurrent.futures.ThreadPoolExecutor(max_workers=4) as executor:
        future_to_file = {executor.submit(process_file, file_path): file_path for file_path in files_to_index}
        for future in concurrent.futures.as_completed(future_to_file):
            file_path = future_to_file[future]
            try:
                doc = future.result()
                if doc:
                    docs_to_index.append(doc)
                    indexed_files[0] += 1
                elif corruptzip.handle_corrupt_file(file_path, corrupt_dir):
                    pass
            except Exception as e:
                logger.error(f"Failed to process {file_path}: {e}")
            
            progress = (indexed_files[0] / total_files * 100) if total_files > 0 else 100
            socketio.emit("index_progress", {
                "progress": progress,
                "indexed": indexed_files[0],
                "total": total_files,
                "corrupt": corrupt_files[0],
                "repaired": repaired_files[0]
            })
            socketio.sleep(0)
    
    if docs_to_index:
        try:
            from elasticsearch.helpers import bulk
            def generate_actions():
                for doc in docs_to_index:
                    yield {"_index": "knowledge_repo", "_source": doc}
            success, failed = bulk(es, generate_actions(), stats_only=True)
            logger.debug(f"Batch indexed {success} documents, {failed} failed")
            if failed:
                logger.error(f"{failed} documents failed to index")
                socketio.emit("index_complete", {"message": f"Error: {failed} documents failed to index"})
                indexed_files[0] = 0
        except Exception as e:
            logger.error(f"Failed to batch index: {e}")
            socketio.emit("index_complete", {"message": f"Error: Failed to batch index: {e}"})
            indexed_files[0] = 0
    
    corrupt_count = corruptzip.zip_corrupt_files(socketio, corrupt_dir, corrupt_zip_path)
    corrupt_files[0] = corrupt_count
    
    es.indices.refresh(index="knowledge_repo")
    message = f"Indexing completed! Indexed {indexed_files[0]} out of {total_files} files."
    if repaired_files[0] > 0:
        message += f" Repaired and indexed {repaired_files[0]} corrupt files."
    if corrupt_files[0] > 0:
        message += f" Skipped and zipped {corrupt_files[0]} unrepairable files to {corrupt_zip_path}."
    logger.debug(message)
    socketio.emit("index_complete", {"message": message})
    return total_files, indexed_files[0]

def search(query, last_indexed_dir=None):
    body = {"query": {"bool": {"must": []}}, "size": 50}
    
    if ":" in query:
        field, value = query.split(":", 1)
        if field.startswith("metadata."):
            body["query"]["bool"]["must"].append({"match": {field: {"query": value, "fuzziness": "AUTO"}}})
        else:
            body["query"]["bool"]["must"].append({"multi_match": {"query": query, "fields": ["content^2", "path"], "fuzziness": "AUTO"}})
    else:
        body["query"]["bool"]["must"].append({"multi_match": {"query": query, "fields": ["content^2", "path"], "fuzziness": "AUTO"}})
    
    if last_indexed_dir:
        body["query"]["bool"]["filter"] = {"prefix": {"path": last_indexed_dir}}
    
    try:
        response = es.search(index="knowledge_repo", body=body)
        hits = response["hits"]["hits"]
        results = []
        if not hits:
            return results
        
        max_score = max(hit["_score"] for hit in hits) if hits else 1
        
        for hit in hits:
            source = hit["_source"]
            content = source["content"]
            snippet = content[:200] + "..." if content else "No content available"
            result = {
                "path": source["path"],
                "snippet": snippet,
                "score": hit["_score"],
                "match_percentage": int((hit["_score"] / max_score) * 100) if max_score > 0 else 0
            }
            if "metadata" in source:
                result["metadata"] = source["metadata"]
            results.append(result)
        results.sort(key=lambda x: x["score"], reverse=True)
        return results
    except Exception as e:
        logger.error(f"Search error: {e}")
        return [{"error": str(e)}]

def libreoffice_to_pdf(file_path, output_path):
    try:
        cmd = [LIBREOFFICE_PATH, "--headless", "--convert-to", "pdf", "--outdir", os.path.dirname(output_path), file_path]
        result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, timeout=30)
        if result.returncode != 0:
            logger.error(f"LibreOffice conversion failed: {result.stderr}")
            return None
        generated_pdf = os.path.join(os.path.dirname(output_path), os.path.splitext(os.path.basename(file_path))[0] + ".pdf")
        if os.path.exists(generated_pdf):
            os.rename(generated_pdf, output_path)
            logger.debug(f"Converted to PDF with LibreOffice: {output_path}")
            return output_path
        else:
            logger.error(f"LibreOffice output not found: {generated_pdf}")
            return None
    except subprocess.TimeoutExpired:
        logger.error(f"LibreOffice conversion timed out for {file_path}")
        return None
    except Exception as e:
        logger.error(f"Error converting {file_path} to PDF with LibreOffice: {e}")
        return None

def docx_to_pdf_fallback(file_path, output_path):
    try:
        doc = Document(file_path)
        temp_pdf = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
        doc_pdf = SimpleDocTemplate(temp_pdf.name, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []

        bold_style = ParagraphStyle('Bold', parent=styles['Normal'], fontName='Helvetica-Bold')
        italic_style = ParagraphStyle('Italic', parent=styles['Normal'], fontName='Helvetica-Oblique')

        for para in doc.paragraphs:
            text = ""
            for run in para.runs:
                if run.bold and run.italic:
                    text += f"<b><i>{run.text}</i></b>"
                elif run.bold:
                    text += f"<b>{run.text}</b>"
                elif run.italic:
                    text += f"<i>{run.text}</i>"
                else:
                    text += run.text
            if text.strip():
                story.append(Paragraph(text, styles["Normal"]))
                story.append(Spacer(1, 12))
        
        doc_pdf.build(story)
        logger.debug(f"Created formatted fallback PDF: {temp_pdf.name}")
        return temp_pdf.name
    except Exception as e:
        logger.error(f"Error creating formatted fallback PDF for {file_path}: {e}")
        return None

@app.route("/preview/<path:file_path>")
def preview_file(file_path):
    try:
        last_indexed_dir = session.get('last_indexed_dir', '/Volumes/CORRUPTED/AI/Sample')
        logger.debug(f"Last indexed directory: {last_indexed_dir}")

        file_path = os.path.normpath(file_path)
        if '..' in file_path or file_path.startswith('/'):
            logger.error(f"Invalid path: {file_path}")
            return jsonify({"error": "Invalid path"}), 400

        full_path = os.path.normpath(os.path.join(last_indexed_dir, file_path.replace('/', os.sep)))
        logger.debug(f"Requested path: {full_path}")

        if not os.path.exists(full_path):
            possible_paths = glob.glob(os.path.join(last_indexed_dir, '*', '*'))
            matching_path = next((p for p in possible_paths if os.path.basename(p).lower() == os.path.basename(full_path).lower()), None)
            if matching_path:
                full_path = matching_path
                logger.debug(f"Matched path: {full_path}")
            else:
                logger.error(f"File not found: {full_path}")
                return jsonify({"error": "File not found"}), 404
        elif not os.path.isfile(full_path):
            logger.error(f"Not a file: {full_path}")
            return jsonify({"error": "Not a file"}), 400

        ext = os.path.splitext(full_path)[1].lower()
        pdf_path = full_path

        if ext == '.docx':
            temp_pdf = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
            try:
                docx_to_pdf(full_path, temp_pdf.name)
                if os.path.getsize(temp_pdf.name) == 0:
                    pdf_path = docx_to_pdf_fallback(full_path, temp_pdf.name) or pdf_path
                else:
                    pdf_path = temp_pdf.name
            except Exception as e:
                logger.error(f"Docx to PDF error: {e}")
                pdf_path = docx_to_pdf_fallback(full_path, temp_pdf.name) or pdf_path
            if not pdf_path:
                return jsonify({"error": "Failed to convert DOCX to PDF"}), 500
        elif ext == '.pptx':
            temp_pdf = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
            pdf_path = libreoffice_to_pdf(full_path, temp_pdf.name)
            if not pdf_path:
                text = extract_text(full_path, skip_ocr=True)
                if text:
                    doc = SimpleDocTemplate(temp_pdf.name, pagesize=letter)
                    styles = getSampleStyleSheet()
                    story = [Paragraph(text.replace('\n', '<br />'), styles["Normal"])]
                    doc.build(story)
                    pdf_path = temp_pdf.name
                else:
                    return jsonify({"error": "Failed to convert PPTX to PDF"}), 500
        elif ext == '.xml':
            temp_pdf = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
            content, metadata = extract_text(full_path, skip_ocr=True)
            doc = SimpleDocTemplate(temp_pdf.name, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []
            if metadata:
                story.append(Paragraph("Metadata:", styles["Heading2"]))
                for key, value in metadata.items():
                    story.append(Paragraph(f"{key}: {value}", styles["Normal"]))
                story.append(Spacer(1, 12))
            if content:
                story.append(Paragraph("Content:", styles["Heading2"]))
                story.append(Paragraph(content.replace('\n', '<br />'), styles["Normal"]))
            doc.build(story)
            pdf_path = temp_pdf.name
            logger.debug(f"Generated XML PDF preview: {pdf_path}")
        elif ext in {".png", ".jpeg", ".jpg"}:
            temp_pdf = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
            content = extract_text(full_path, skip_ocr=False)
            doc = SimpleDocTemplate(temp_pdf.name, pagesize=letter)
            styles = getSampleStyleSheet()
            story = [Paragraph("Extracted Text from Image:", styles["Heading2"])]
            story.append(Paragraph(content.replace('\n', '<br />'), styles["Normal"]))
            doc.build(story)
            pdf_path = temp_pdf.name
            logger.debug(f"Generated image PDF preview: {pdf_path}")
        elif ext in {".xlsx", ".csv"}:
            temp_pdf = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
            content = extract_text(full_path, skip_ocr=True)
            doc = SimpleDocTemplate(temp_pdf.name, pagesize=letter)
            styles = getSampleStyleSheet()
            story = [Paragraph("Content:", styles["Heading2"])]
            story.append(Paragraph(content.replace('\n', '<br />'), styles["Normal"]))
            doc.build(story)
            pdf_path = temp_pdf.name
            logger.debug(f"Generated {ext} PDF preview: {pdf_path}")
        elif ext != '.pdf':
            logger.error(f"Unsupported preview format: {ext}")
            return jsonify({"error": "Preview not supported for this file type"}), 400

        with open(pdf_path, 'rb') as f:
            file_data = f.read()
        file_size = len(file_data)
        if file_size == 0:
            logger.error(f"PDF empty: {pdf_path}")
            return jsonify({"error": "PDF empty"}), 500
        response = make_response(file_data)
        response.headers['Content-Type'] = 'application/pdf'
        response.headers['Content-Disposition'] = f'inline; filename="{os.path.basename(full_path)}.pdf"'
        response.headers['Content-Length'] = file_size
        response.headers['Access-Control-Allow-Origin'] = '*'
        if pdf_path != full_path:
            os.unlink(pdf_path)
        return response
    except Exception as e:
        logger.error(f"Preview error: {file_path} - {e}", exc_info=True)
        return jsonify({"error": str(e)}), 500

@app.route('/index_directory', methods=['POST'])
def index_directory_route():
    directory = request.form.get('directory', '/Volumes/CORRUPTED/AI/Sample')
    logger.debug(f"Indexing directory: {directory}")
    if not os.path.isdir(directory):
        return jsonify({"error": f"Directory not found: {directory}"}), 400
    session['last_indexed_dir'] = directory
    logger.debug(f"Set last_indexed_dir to: {session['last_indexed_dir']}")
    files = [f for f in os.listdir(directory) if os.path.isfile(os.path.join(directory, f))]
    return jsonify({"status": "success", "directory": directory, "files": files})

@app.route("/serve/<path:file_path>")
def serve_file(file_path):
    try:
        last_indexed_dir = session.get('last_indexed_dir', '/Volumes/CORRUPTED/AI/Sample')
        logger.debug(f"Session last_indexed_dir: {last_indexed_dir}")
        file_path = urllib.parse.unquote(file_path).replace('/', os.sep)
        full_path = os.path.normpath(os.path.join(last_indexed_dir, file_path))
        logger.debug(f"Attempting to serve file: {full_path}")

        if not os.path.exists(full_path):
            logger.error(f"File not found: {full_path}. Directory contents: {os.listdir(last_indexed_dir)}")
            return jsonify({"error": f"File not found: {full_path}", "dir_contents": os.listdir(last_indexed_dir)}), 404
        elif not os.path.isfile(full_path):
            logger.error(f"Path is not a file: {full_path}")
            return jsonify({"error": "Path is not a file"}), 400

        if not os.access(full_path, os.R_OK):
            logger.error(f"File not readable: {full_path}")
            return jsonify({"error": "File not readable"}), 403

        ext = os.path.splitext(full_path)[1].lower()
        content_types = {
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            '.doc': 'application/msword',
            '.ppt': 'application/vnd.ms-powerpoint',
            '.xls': 'application/vnd.ms-excel',
            '.pdf': 'application/pdf',
            '.png': 'image/png',
            '.jpeg': 'image/jpeg',
            '.jpg': 'image/jpeg',
            '.xml': 'application/xml'
        }
        content_type = content_types.get(ext, 'application/octet-stream')

        with open(full_path, 'rb') as f:
            file_data = f.read()
        file_size = len(file_data)

        if file_size == 0 or file_size > 25 * 1024 * 1024:
            logger.error(f"File empty or >25MB: {full_path} (size: {file_size} bytes)")
            return jsonify({"error": "File empty or exceeds 25MB limit"}), 400

        response = make_response(file_data)
        response.headers['Content-Type'] = content_type
        response.headers['Content-Length'] = file_size
        response.headers['Content-Disposition'] = 'inline'
        response.headers['Access-Control-Allow-Origin'] = '*'
        response.headers['Access-Control-Allow-Methods'] = 'GET, OPTIONS'
        response.headers['Access-Control-Allow-Headers'] = 'Content-Type'
        response.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate'
        response.headers['Vary'] = 'Origin'
        logger.debug(f"Served file: {full_path} with Content-Type: {content_type}")
        return response
    except PermissionError as e:
        logger.error(f"Permission denied: {full_path} - {e}")
        return jsonify({"error": "Permission denied"}), 403
    except Exception as e:
        logger.error(f"Error serving {file_path}: {e}", exc_info=True)
        return jsonify({"error": str(e)}), 500

@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        if "index" in request.form:
            directory = request.form.get("directory", "").strip()
            if not directory:
                return render_template("index.html", message="Please enter or select a directory path.")
            directory = os.path.expanduser(directory)
            if os.path.isdir(directory):
                session['last_indexed_dir'] = directory
                return render_template("index.html", indexing=True, directory=directory, last_indexed_dir=directory, public_url=app.config['public_url'])
            return render_template("index.html", message=f"Invalid directory path: {directory}")
        elif "search" in request.form:
            query = request.form.get("query", "").strip()
            if query:
                last_indexed_dir = session.get('last_indexed_dir')
                results = search(query, last_indexed_dir)
                return render_template("index.html", results=results, query=query, count=len(results), last_indexed_dir=last_indexed_dir, public_url=app.config['public_url'])
            return render_template("index.html", message="Please enter a search query.")
    return render_template("index.html", last_indexed_dir=session.get('last_indexed_dir', '/Volumes/CORRUPTED/AI/Sample'), public_url=app.config['public_url'])

@socketio.on("start_indexing")
def handle_start_indexing(data):
    logger.debug("Received start_indexing event with data: %s", data)
    directory = data.get("directory", "").strip()
    if not directory:
        emit("index_complete", {"message": "Error: No directory provided."})
        return
    directory = os.path.expanduser(directory)
    if os.path.isdir(directory):
        session['last_indexed_dir'] = directory
        total_files, indexed_files = index_directory(directory)
        emit("index_complete", {"message": f"Indexing completed! Indexed {indexed_files} out of {total_files} files."})
    else:
        emit("index_complete", {"message": f"Error: Directory {directory} is not valid or accessible."})

@app.route("/update_public_url", methods=["POST"])
def update_public_url():
    new_url = request.form.get("public_url", "").strip()
    if new_url:
        app.config['public_url'] = new_url
        session['public_url'] = new_url
        logger.debug(f"Updated public URL to: {new_url}")
        return jsonify({"status": "success", "public_url": new_url})
    return jsonify({"status": "error", "message": "Invalid public URL"}), 400

if __name__ == "__main__":
    if not es.ping():
        logger.error("Elasticsearch is not running. Start it with 'brew services start elastic/tap/elasticsearch-full'")
        exit(1)
    socketio.run(app, debug=True)